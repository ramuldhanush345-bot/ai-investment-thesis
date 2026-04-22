from flask import Flask, request
from pptx import Presentation

app = Flask(__name__)

@app.route('/extract', methods=['POST'])
def extract():
    if 'file' not in request.files:
        return "No file"

    file = request.files['file']
    prs = Presentation(file)

    all_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        all_text.append(paragraph.text.strip())

    return " ".join(all_text)


if __name__ == "__main__":
    print("Python service running")
    app.run(port=7000)